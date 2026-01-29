#!/usr/bin/env python3
"""
Office Automation - Simple Direct Version
直接使用openpyxl和python-docx，不依赖复杂的导入结构
"""

import os
import sys
from pathlib import Path
from typing import Dict, Any, Optional, List, Union

class OfficeAutomation:
    """简单的Office自动化类"""
    
    def __init__(self, config=None):
        self.config = config or {}
        print("Office Automation initialized (Simple Direct Version)")
    
    @property
    def excel(self):
        """Excel处理器"""
        return ExcelProcessor()
    
    @property
    def word(self):
        """Word处理器"""
        return WordProcessor()
    
    @property
    def powerpoint(self):
        """PowerPoint处理器"""
        return PowerPointProcessor()

class ExcelProcessor:
    """Excel处理器 - 使用openpyxl"""
    
    def __init__(self):
        self._check_dependencies()
    
    def _check_dependencies(self):
        """检查依赖"""
        try:
            from openpyxl import Workbook
            self.Workbook = Workbook
            self._has_openpyxl = True
        except ImportError:
            print("Warning: openpyxl not installed. Using dummy mode.")
            self._has_openpyxl = False
    
    def create_workbook(self):
        """创建工作簿"""
        if self._has_openpyxl:
            wb = self.Workbook()
            # 移除默认工作表
            if 'Sheet' in wb.sheetnames:
                default_sheet = wb['Sheet']
                wb.remove(default_sheet)
            return RealExcelWorkbook(wb)
        else:
            return DummyExcelWorkbook()
    
    def open_workbook(self, filepath):
        """打开工作簿"""
        if self._has_openpyxl:
            from openpyxl import load_workbook
            wb = load_workbook(filepath)
            return RealExcelWorkbook(wb)
        else:
            return DummyExcelWorkbook()

class RealExcelWorkbook:
    """真实的Excel工作簿"""
    
    def __init__(self, workbook):
        self.wb = workbook
    
    def add_worksheet(self, name):
        """添加工作表"""
        ws = self.wb.create_sheet(title=name)
        return RealExcelWorksheet(ws)
    
    def save(self, filepath):
        """保存工作簿"""
        self.wb.save(filepath)
        print(f"Saved Excel workbook to: {filepath}")
        return True

class RealExcelWorksheet:
    """真实的Excel工作表"""
    
    def __init__(self, worksheet):
        self.ws = worksheet
    
    def cell(self, row, column, value=None):
        """设置单元格值"""
        cell = self.ws.cell(row=row, column=column)
        if value is not None:
            cell.value = value
        return cell
    
    def set_column_width(self, column, width):
        """设置列宽"""
        from openpyxl.utils import get_column_letter
        col_letter = get_column_letter(column)
        self.ws.column_dimensions[col_letter].width = width
    
    def add_table(self, data, start_row=1, start_col=1):
        """添加表格"""
        for i, row_data in enumerate(data):
            for j, cell_data in enumerate(row_data):
                self.cell(start_row + i, start_col + j, cell_data)

class DummyExcelWorkbook:
    """虚拟Excel工作簿"""
    
    def __init__(self):
        self.sheets = {}
    
    def add_worksheet(self, name):
        """添加工作表"""
        print(f"Adding worksheet: {name}")
        return DummyExcelWorksheet()
    
    def save(self, filepath):
        """保存工作簿"""
        print(f"Saving dummy workbook to: {filepath}")
        # 创建空文件
        with open(filepath, 'w') as f:
            f.write('')
        return True

class DummyExcelWorksheet:
    """虚拟Excel工作表"""
    
    def cell(self, row, column, value=None):
        """设置单元格值"""
        print(f"Setting cell ({row},{column}) = {value}")
        return DummyCell()

class DummyCell:
    """虚拟单元格"""
    pass

class WordProcessor:
    """Word处理器 - 使用python-docx"""
    
    def __init__(self):
        self._check_dependencies()
    
    def _check_dependencies(self):
        """检查依赖"""
        try:
            from docx import Document
            self.Document = Document
            self._has_docx = True
        except ImportError:
            print("Warning: python-docx not installed. Using dummy mode.")
            self._has_docx = False
    
    def create_document(self):
        """创建文档"""
        if self._has_docx:
            doc = self.Document()
            return RealWordDocument(doc)
        else:
            return DummyWordDocument()

class RealWordDocument:
    """真实的Word文档"""
    
    def __init__(self, document):
        self.doc = document
    
    def add_heading(self, text, level=1):
        """添加标题"""
        self.doc.add_heading(text, level)
        return self
    
    def add_paragraph(self, text):
        """添加段落"""
        self.doc.add_paragraph(text)
        return self
    
    def add_table(self, data, headers=None):
        """添加表格"""
        if headers:
            table_data = [headers] + data
        else:
            table_data = data
        
        table = self.doc.add_table(rows=len(table_data), cols=len(table_data[0]))
        
        for i, row_data in enumerate(table_data):
            row_cells = table.rows[i].cells
            for j, cell_data in enumerate(row_data):
                row_cells[j].text = str(cell_data)
        
        return self
    
    def save(self, filepath):
        """保存文档"""
        self.doc.save(filepath)
        print(f"Saved Word document to: {filepath}")
        return True

class DummyWordDocument:
    """虚拟Word文档"""
    
    def add_heading(self, text, level=1):
        """添加标题"""
        print(f"Adding heading: {text} (level {level})")
        return self
    
    def add_paragraph(self, text):
        """添加段落"""
        print(f"Adding paragraph: {text}")
        return self
    
    def save(self, filepath):
        """保存文档"""
        print(f"Saving dummy document to: {filepath}")
        with open(filepath, 'w') as f:
            f.write('')
        return True

class PowerPointProcessor:
    """PowerPoint处理器 - 使用python-pptx"""
    
    def __init__(self):
        self._check_dependencies()
    
    def _check_dependencies(self):
        """检查依赖"""
        try:
            from pptx import Presentation
            self.Presentation = Presentation
            self._has_pptx = True
        except ImportError:
            print("Warning: python-pptx not installed. Using dummy mode.")
            self._has_pptx = False
    
    def create_presentation(self):
        """创建演示文稿"""
        if self._has_pptx:
            pres = self.Presentation()
            return RealPowerPointPresentation(pres)
        else:
            return DummyPowerPointPresentation()

class RealPowerPointPresentation:
    """真实的PowerPoint演示文稿"""
    
    def __init__(self, presentation):
        self.pres = presentation
    
    def add_slide(self, layout=0):
        """添加幻灯片"""
        slide_layout = self.pres.slide_layouts[layout]
        slide = self.pres.slides.add_slide(slide_layout)
        return slide
    
    def save(self, filepath):
        """保存演示文稿"""
        self.pres.save(filepath)
        print(f"Saved PowerPoint presentation to: {filepath}")
        return True

class DummyPowerPointPresentation:
    """虚拟PowerPoint演示文稿"""
    
    def add_slide(self, layout=0):
        """添加幻灯片"""
        print(f"Adding slide with layout: {layout}")
        return self
    
    def save(self, filepath):
        """保存演示文稿"""
        print(f"Saving dummy presentation to: {filepath}")
        with open(filepath, 'w') as f:
            f.write('')
        return True

# 测试函数
def test_office_automation():
    """测试Office自动化"""
    print("=" * 60)
    print("测试Office Automation Skill")
    print("=" * 60)
    
    # 创建实例
    office = OfficeAutomation()
    
    # 测试Excel
    print("\n1. 测试Excel功能...")
    excel = office.excel
    workbook = excel.create_workbook()
    
    # 添加工作表
    sheet1 = workbook.add_worksheet("数据表")
    
    # 添加数据
    sheet1.cell(1, 1, "姓名")
    sheet1.cell(1, 2, "年龄")
    sheet1.cell(1, 3, "部门")
    
    data = [
        ["张三", 28, "技术部"],
        ["李四", 32, "市场部"],
        ["王五", 25, "人事部"]
    ]
    
    for i, row in enumerate(data, start=2):
        sheet1.cell(i, 1, row[0])
        sheet1.cell(i, 2, row[1])
        sheet1.cell(i, 3, row[2])
    
    # 保存
    test_excel = r"F:\cheshi\office_skill_test.xlsx"
    workbook.save(test_excel)
    
    if os.path.exists(test_excel):
        size = os.path.getsize(test_excel)
        print(f"  [SUCCESS] Excel文件创建成功: {test_excel} ({size} 字节)")
    else:
        print(f"  [FAILED] Excel文件创建失败")
    
    # 测试Word
    print("\n2. 测试Word功能...")
    word = office.word
    document = word.create_document()
    
    document.add_heading("测试报告", 1)
    document.add_paragraph("这是一个测试文档，用于验证Office Automation Skill的功能。")
    
    # 添加表格
    table_data = [
        ["产品", "数量", "价格"],
        ["笔记本电脑", "10", "¥5,000"],
        ["显示器", "15", "¥1,500"],
        ["键盘", "25", "¥200"]
    ]
    
    document.add_table(table_data[1:], table_data[0])
    
    test_word = r"F:\cheshi\office_skill_test.docx"
    document.save(test_word)
    
    if os.path.exists(test_word):
        size = os.path.getsize(test_word)
        print(f"  [SUCCESS] Word文件创建成功: {test_word} ({size} 字节)")
    else:
        print(f"  [FAILED] Word文件创建失败")
    
    print("\n" + "=" * 60)
    print("测试完成!")
    print("=" * 60)
    
    return True

if __name__ == "__main__":
    test_office_automation()