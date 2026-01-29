"""
Format Converter for Office Documents

This module provides functionality for converting between different
Office document formats with WPS compatibility.
"""

import os
from typing import Optional, List, Dict, Any, Union
from pathlib import Path
import warnings

from ..utils.error_handler import (
    FormatConversionError,
    ValidationError,
    create_error_context,
    validate_file_path,
    validate_format,
)


class FormatConverter:
    """Converter for Office document formats."""
    
    # Supported conversion formats
    SUPPORTED_CONVERSIONS = {
        # Word conversions
        "docx": ["pdf", "txt", "rtf", "html"],
        "doc": ["docx", "pdf", "txt"],
        
        # Excel conversions  
        "xlsx": ["csv", "pdf", "html", "json"],
        "xls": ["xlsx", "csv", "pdf"],
        
        # PowerPoint conversions
        "pptx": ["pdf", "jpg", "png"],
        "ppt": ["pptx", "pdf"],
        
        # PDF conversions
        "pdf": ["docx", "txt", "jpg", "png"],
    }
    
    def __init__(self):
        """Initialize Format Converter."""
        self._check_dependencies()
    
    def _check_dependencies(self) -> None:
        """Check if required dependencies are available."""
        # These are optional dependencies for advanced conversions
        # The basic converters will work with the core libraries
        
        optional_deps = {
            "pdf_conversion": ["pdf2docx", "PyPDF2", "reportlab"],
            "image_conversion": ["PIL", "pillow"],
            "advanced_conversion": ["unoconv", "libreoffice"],
        }
        
        missing_deps = []
        for feature, deps in optional_deps.items():
            for dep in deps:
                try:
                    __import__(dep.replace("-", "_"))
                except ImportError:
                    missing_deps.append((feature, dep))
        
        if missing_deps:
            warnings.warn(
                f"Some optional dependencies missing for advanced features: {missing_deps}\n"
                "Basic conversions will work, but some features may be limited."
            )
    
    def convert(
        self,
        input_file: Union[str, Path],
        output_format: str,
        output_file: Optional[Union[str, Path]] = None,
        options: Optional[Dict[str, Any]] = None
    ) -> str:
        """
        Convert a document from one format to another.
        
        Args:
            input_file: Path to input file
            output_format: Desired output format
            output_file: Path for output file (optional, auto-generated if None)
            options: Conversion options
            
        Returns:
            Path to converted file
            
        Raises:
            FormatConversionError: If conversion fails
        """
        context = create_error_context(
            "convert_format",
            file_path=input_file,
            format=output_format
        )
        
        input_file = Path(input_file)
        validate_file_path(str(input_file), "convert")
        
        # Auto-generate output filename if not provided
        if output_file is None:
            output_file = input_file.with_suffix(f".{output_format}")
        else:
            output_file = Path(output_file)
        
        validate_file_path(str(output_file), "save")
        
        # Get input format from file extension
        input_format = input_file.suffix.lower().lstrip('.')
        
        # Validate conversion is supported
        self._validate_conversion(input_format, output_format)
        
        try:
            # Perform conversion based on formats
            if input_format == "docx" and output_format == "pdf":
                return self._docx_to_pdf(input_file, output_file, options)
            elif input_format == "xlsx" and output_format == "csv":
                return self._xlsx_to_csv(input_file, output_file, options)
            elif input_format == "pptx" and output_format == "pdf":
                return self._pptx_to_pdf(input_file, output_file, options)
            elif input_format == "pdf" and output_format == "docx":
                return self._pdf_to_docx(input_file, output_file, options)
            else:
                # Generic conversion using fallback method
                return self._generic_conversion(input_file, output_file, options)
                
        except Exception as e:
            raise FormatConversionError(
                f"Failed to convert {input_format} to {output_format}: {e}",
                context
            )
    
    def batch_convert(
        self,
        input_files: List[Union[str, Path]],
        output_format: str,
        output_dir: Optional[Union[str, Path]] = None,
        options: Optional[Dict[str, Any]] = None
    ) -> List[str]:
        """
        Convert multiple files to the specified format.
        
        Args:
            input_files: List of input file paths
            output_format: Desired output format
            output_dir: Output directory (optional, uses input directory if None)
            options: Conversion options
            
        Returns:
            List of paths to converted files
            
        Raises:
            FormatConversionError: If any conversion fails
        """
        results = []
        errors = []
        
        for input_file in input_files:
            try:
                # Determine output directory
                if output_dir:
                    output_dir_path = Path(output_dir)
                    output_dir_path.mkdir(parents=True, exist_ok=True)
                    input_path = Path(input_file)
                    output_file = output_dir_path / f"{input_path.stem}.{output_format}"
                else:
                    output_file = None
                
                # Convert file
                result = self.convert(input_file, output_format, output_file, options)
                results.append(result)
                
            except Exception as e:
                errors.append(f"{input_file}: {e}")
        
        if errors:
            raise FormatConversionError(
                f"Batch conversion completed with errors:\n" + "\n".join(errors)
            )
        
        return results
    
    def _validate_conversion(self, input_format: str, output_format: str) -> None:
        """
        Validate that the requested conversion is supported.
        
        Args:
            input_format: Input file format
            output_format: Output file format
            
        Raises:
            FormatConversionError: If conversion is not supported
        """
        if input_format not in self.SUPPORTED_CONVERSIONS:
            raise FormatConversionError(
                f"Input format '{input_format}' is not supported. "
                f"Supported formats: {list(self.SUPPORTED_CONVERSIONS.keys())}"
            )
        
        if output_format not in self.SUPPORTED_CONVERSIONS[input_format]:
            raise FormatConversionError(
                f"Conversion from '{input_format}' to '{output_format}' is not supported. "
                f"Supported conversions from '{input_format}': {self.SUPPORTED_CONVERSIONS[input_format]}"
            )
    
    def _docx_to_pdf(
        self,
        input_file: Path,
        output_file: Path,
        options: Optional[Dict[str, Any]] = None
    ) -> str:
        """Convert DOCX to PDF."""
        try:
            # Try using python-docx with reportlab
            from docx import Document
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            
            # Load Word document
            doc = Document(str(input_file))
            
            # Create PDF
            c = canvas.Canvas(str(output_file), pagesize=letter)
            width, height = letter
            
            # Extract text from Word document
            text_content = []
            for paragraph in doc.paragraphs:
                text_content.append(paragraph.text)
            
            # Simple text rendering (this is simplified)
            # In production, you'd use a proper DOCX to PDF converter
            y = height - 40
            for line in text_content:
                c.drawString(40, y, line[:100])  # Limit line length
                y -= 20
                if y < 40:
                    c.showPage()
                    y = height - 40
            
            c.save()
            
            warnings.warn(
                "DOCX to PDF conversion is simplified. "
                "For production use, install a proper converter like 'docx2pdf'."
            )
            
            return str(output_file)
            
        except ImportError:
            # Fallback: Save as text and inform user
            raise FormatConversionError(
                "PDF conversion requires additional libraries. "
                "Install: pip install reportlab"
            )
    
    def _xlsx_to_csv(
        self,
        input_file: Path,
        output_file: Path,
        options: Optional[Dict[str, Any]] = None
    ) -> str:
        """Convert XLSX to CSV."""
        try:
            import openpyxl
            import csv
            
            # Load workbook
            wb = openpyxl.load_workbook(str(input_file), data_only=True)
            
            # Get active worksheet or first worksheet
            if options and "sheet_name" in options:
                ws = wb[options["sheet_name"]]
            else:
                ws = wb.active
            
            # Write to CSV
            with open(output_file, 'w', newline='', encoding='utf-8') as f:
                writer = csv.writer(f)
                
                # Write all rows
                for row in ws.iter_rows(values_only=True):
                    writer.writerow(row)
            
            return str(output_file)
            
        except Exception as e:
            raise FormatConversionError(f"Failed to convert XLSX to CSV: {e}")
    
    def _pptx_to_pdf(
        self,
        input_file: Path,
        output_file: Path,
        options: Optional[Dict[str, Any]] = None
    ) -> str:
        """Convert PPTX to PDF."""
        # This is a placeholder - python-pptx doesn't have built-in PDF export
        raise FormatConversionError(
            "PPTX to PDF conversion requires additional setup. "
            "Options:\n"
            "1. Use Microsoft Office or LibreOffice in headless mode\n"
            "2. Use a cloud conversion service\n"
            "3. Install a commercial PPTX to PDF converter"
        )
    
    def _pdf_to_docx(
        self,
        input_file: Path,
        output_file: Path,
        options: Optional[Dict[str, Any]] = None
    ) -> str:
        """Convert PDF to DOCX."""
        try:
            # Try using pdf2docx
            from pdf2docx import Converter
            
            cv = Converter(str(input_file))
            cv.convert(str(output_file))
            cv.close()
            
            return str(output_file)
            
        except ImportError:
            raise FormatConversionError(
                "PDF to DOCX conversion requires 'pdf2docx'. "
                "Install: pip install pdf2docx"
            )
        except Exception as e:
            raise FormatConversionError(f"Failed to convert PDF to DOCX: {e}")
    
    def _generic_conversion(
        self,
        input_file: Path,
        output_file: Path,
        options: Optional[Dict[str, Any]] = None
    ) -> str:
        """
        Generic conversion method for unsupported conversions.
        
        This method provides a fallback by extracting text content
        and saving in the new format.
        """
        # Extract text based on file type
        text_content = self._extract_text(input_file)
        
        # Save in new format
        if output_file.suffix.lower() == ".txt":
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(text_content)
        elif output_file.suffix.lower() == ".html":
            html_content = f"<html><body><pre>{text_content}</pre></body></html>"
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(html_content)
        else:
            # For other formats, save as text with warning
            warnings.warn(
                f"Conversion to {output_file.suffix} is simplified. "
                "Only text content will be preserved."
            )
            with open(output_file.with_suffix('.txt'), 'w', encoding='utf-8') as f:
                f.write(text_content)
            output_file = output_file.with_suffix('.txt')
        
        return str(output_file)
    
    def _extract_text(self, file_path: Path) -> str:
        """Extract text content from various file formats."""
        suffix = file_path.suffix.lower()
        
        try:
            if suffix == ".docx":
                from docx import Document
                doc = Document(str(file_path))
                return "\n".join([p.text for p in doc.paragraphs])
            
            elif suffix == ".xlsx":
                import openpyxl
                wb = openpyxl.load_workbook(str(file_path), data_only=True)
                ws = wb.active
                text_lines = []
                for row in ws.iter_rows(values_only=True):
                    text_lines.append(",".join([str(cell) if cell is not None else "" for cell in row]))
                return "\n".join(text_lines)
            
            elif suffix == ".pptx":
                from pptx import Presentation
                prs = Presentation(str(file_path))
                text_lines = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_lines.append(shape.text)
                return "\n".join(text_lines)
            
            elif suffix == ".pdf":
                try:
                    import PyPDF2
                    with open(file_path, 'rb') as f:
                        reader = PyPDF2.PdfReader(f)
                        text = ""
                        for page in reader.pages:
                            text += page.extract_text()
                        return text
                except ImportError:
                    return f"[PDF content extraction requires PyPDF2: {file_path}]"
            
            elif suffix == ".txt":
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            
            else:
                return f"[Unsupported file format for text extraction: {suffix}]"
                
        except Exception as e:
            return f"[Error extracting text from {file_path}: {e}]"
    
    def get_supported_conversions(self) -> Dict[str, List[str]]:
        """Get dictionary of supported format conversions."""
        return self.SUPPORTED_CONVERSIONS.copy()
    
    def is_conversion_supported(self, input_format: str, output_format: str) -> bool:
        """Check if a conversion is supported."""
        if input_format not in self.SUPPORTED_CONVERSIONS:
            return False
        return output_format in self.SUPPORTED_CONVERSIONS[input_format]
    
    def add_conversion_support(
        self,
        input_format: str,
        output_formats: List[str]
    ) -> None:
        """
        Add support for new format conversions.
        
        Args:
            input_format: Input file format
            output_formats: List of supported output formats
        """
        if input_format in self.SUPPORTED_CONVERSIONS:
            self.SUPPORTED_CONVERSIONS[input_format].extend(output_formats)
            # Remove duplicates
            self.SUPPORTED_CONVERSIONS[input_format] = list(
                set(self.SUPPORTED_CONVERSIONS[input_format])
            )
        else:
            self.SUPPORTED_CONVERSIONS[input_format] = output_formats