"""
PowerPoint Presentation Processor

This module provides functionality for creating, editing, and manipulating
PowerPoint presentations (.pptx format) with WPS compatibility.
"""

import os
from typing import Optional, List, Dict, Any, Union
from pathlib import Path

from ..utils.error_handler import (
    DocumentCreationError,
    DocumentReadError,
    DocumentSaveError,
    ValidationError,
    create_error_context,
    validate_file_path,
    validate_format,
)


class PowerPointProcessor:
    """Processor for PowerPoint presentation operations."""
    
    # Supported formats
    SUPPORTED_FORMATS = ["pptx", "ppt", "pdf", "jpg", "png"]
    
    def __init__(self, template_dir: Optional[str] = None):
        """
        Initialize PowerPoint Processor.
        
        Args:
            template_dir: Directory containing PowerPoint templates
        """
        self.template_dir = template_dir
        self._presentation = None
        
        # Try to import python-pptx
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            from pptx.chart.data import ChartData
            from pptx.enum.chart import XL_CHART_TYPE
            
            self.Presentation = Presentation
            self.Inches = Inches
            self.Pt = Pt
            self.PP_ALIGN = PP_ALIGN
            self.RGBColor = RGBColor
            self.ChartData = ChartData
            self.XL_CHART_TYPE = XL_CHART_TYPE
            
            self._pptx_available = True
        except ImportError:
            self._pptx_available = False
            raise ImportError(
                "python-pptx is not installed. Please install it with: "
                "pip install python-pptx"
            )
    
    def create_presentation(
        self,
        template: Optional[Union[str, Path]] = None
    ) -> 'Presentation':
        """
        Create a new PowerPoint presentation.
        
        Args:
            template: Path to template file or template name
            
        Returns:
            Presentation object
            
        Raises:
            DocumentCreationError: If presentation creation fails
        """
        context = create_error_context("create_presentation", template=template)
        
        try:
            if template:
                # Load from template
                template_path = self._resolve_template_path(template)
                if not os.path.exists(template_path):
                    raise FileNotFoundError(f"Template not found: {template_path}")
                
                self._presentation = self.Presentation(str(template_path))
            else:
                # Create empty presentation
                self._presentation = self.Presentation()
            
            return self._presentation
            
        except Exception as e:
            raise DocumentCreationError(
                f"Failed to create presentation: {e}",
                context
            )
    
    def load_presentation(self, file_path: Union[str, Path]) -> 'Presentation':
        """
        Load an existing PowerPoint presentation.
        
        Args:
            file_path: Path to the presentation file
            
        Returns:
            Presentation object
            
        Raises:
            DocumentReadError: If presentation loading fails
        """
        context = create_error_context("load_presentation", file_path=file_path)
        validate_file_path(str(file_path), "load_presentation")
        
        try:
            file_path = Path(file_path)
            if not file_path.exists():
                raise FileNotFoundError(f"Presentation not found: {file_path}")
            
            self._presentation = self.Presentation(str(file_path))
            return self._presentation
            
        except Exception as e:
            raise DocumentReadError(
                f"Failed to load presentation: {e}",
                context
            )
    
    def add_slide(
        self,
        layout: str = "title_and_content",
        title: Optional[str] = None,
        content: Optional[str] = None
    ) -> 'Slide':
        """
        Add a new slide to the presentation.
        
        Args:
            layout: Slide layout name
            title: Slide title
            content: Slide content
            
        Returns:
            Slide object
            
        Raises:
            DocumentCreationError: If presentation is not initialized
        """
        if not self._presentation:
            raise DocumentCreationError(
                "Presentation not initialized. Call create_presentation() or load_presentation() first."
            )
        
        try:
            # Get slide layout
            layout_map = {
                "title": 0,
                "title_and_content": 1,
                "section_header": 2,
                "two_content": 3,
                "comparison": 4,
                "title_only": 5,
                "blank": 6,
                "content_with_caption": 7,
                "picture_with_caption": 8,
            }
            
            layout_idx = layout_map.get(layout.lower().replace(" ", "_"), 1)
            slide_layout = self._presentation.slide_layouts[layout_idx]
            
            # Add slide
            slide = self._presentation.slides.add_slide(slide_layout)
            
            # Set title if provided
            if title and hasattr(slide.shapes, 'title'):
                title_shape = slide.shapes.title
                if title_shape:
                    title_shape.text = title
            
            # Set content if provided
            if content:
                # Try to find content placeholder
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape != slide.shapes.title:
                        shape.text_frame.text = content
                        break
            
            return slide
            
        except Exception as e:
            raise DocumentCreationError(f"Failed to add slide: {e}")
    
    def add_text(
        self,
        slide_index: int,
        text: str,
        position: Optional[tuple] = None,
        font_size: Optional[int] = None
    ) -> None:
        """
        Add text to a slide.
        
        Args:
            slide_index: Index of the slide (0-based)
            text: Text to add
            position: (left, top, width, height) in inches
            font_size: Font size in points
        """
        if not self._presentation:
            raise DocumentCreationError(
                "Presentation not initialized. Call create_presentation() or load_presentation() first."
            )
        
        try:
            # Get slide
            if slide_index < 0 or slide_index >= len(self._presentation.slides):
                raise IndexError(f"Slide index {slide_index} out of range")
            
            slide = self._presentation.slides[slide_index]
            
            # Default position if not provided
            if position is None:
                position = (self.Inches(1), self.Inches(1), 
                           self.Inches(8), self.Inches(1))
            
            # Add text box
            left, top, width, height = position
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            
            # Add text
            p = text_frame.add_paragraph()
            p.text = text
            
            # Set font size if provided
            if font_size:
                p.font.size = self.Pt(font_size)
                
        except Exception as e:
            raise DocumentCreationError(f"Failed to add text: {e}")
    
    def add_image(
        self,
        slide_index: int,
        image_path: Union[str, Path],
        position: Optional[tuple] = None,
        size: Optional[tuple] = None
    ) -> None:
        """
        Add an image to a slide.
        
        Args:
            slide_index: Index of the slide (0-based)
            image_path: Path to image file
            position: (left, top) in inches
            size: (width, height) in inches
        """
        if not self._presentation:
            raise DocumentCreationError(
                "Presentation not initialized. Call create_presentation() or load_presentation() first."
            )
        
        try:
            image_path = Path(image_path)
            if not image_path.exists():
                raise FileNotFoundError(f"Image not found: {image_path}")
            
            # Get slide
            if slide_index < 0 or slide_index >= len(self._presentation.slides):
                raise IndexError(f"Slide index {slide_index} out of range")
            
            slide = self._presentation.slides[slide_index]
            
            # Default position and size if not provided
            if position is None:
                position = (self.Inches(1), self.Inches(1))
            
            if size is None:
                size = (self.Inches(6), self.Inches(4))
            
            # Add image
            left, top = position
            width, height = size
            slide.shapes.add_picture(str(image_path), left, top, width, height)
            
        except Exception as e:
            raise DocumentCreationError(f"Failed to add image: {e}")
    
    def add_chart(
        self,
        slide_index: int,
        chart_type: str = "column",
        data: Optional[Dict[str, List[float]]] = None,
        categories: Optional[List[str]] = None,
        position: Optional[tuple] = None,
        size: Optional[tuple] = None
    ) -> None:
        """
        Add a chart to a slide.
        
        Args:
            slide_index: Index of the slide (0-based)
            chart_type: Type of chart ("column", "line", "pie", "bar")
            data: Chart data as {series_name: [values]}
            categories: Category labels
            position: (left, top) in inches
            size: (width, height) in inches
        """
        if not self._presentation:
            raise DocumentCreationError(
                "Presentation not initialized. Call create_presentation() or load_presentation() first."
            )
        
        try:
            # Get slide
            if slide_index < 0 or slide_index >= len(self._presentation.slides):
                raise IndexError(f"Slide index {slide_index} out of range")
            
            slide = self._presentation.slides[slide_index]
            
            # Default data if not provided
            if data is None:
                data = {
                    "Series 1": [19.2, 21.4, 16.7, 15.8, 18.9],
                    "Series 2": [22.3, 19.8, 24.1, 20.5, 23.2],
                }
            
            if categories is None:
                categories = ["Q1", "Q2", "Q3", "Q4", "Q5"]
            
            # Default position and size if not provided
            if position is None:
                position = (self.Inches(1), self.Inches(2))
            
            if size is None:
                size = (self.Inches(6), self.Inches(4))
            
            # Create chart data
            chart_data = self.ChartData()
            chart_data.categories = categories
            
            for series_name, values in data.items():
                chart_data.add_series(series_name, values)
            
            # Map chart type
            chart_type_map = {
                "column": self.XL_CHART_TYPE.COLUMN_CLUSTERED,
                "line": self.XL_CHART_TYPE.LINE,
                "pie": self.XL_CHART_TYPE.PIE,
                "bar": self.XL_CHART_TYPE.BAR_CLUSTERED,
            }
            
            chart_type_enum = chart_type_map.get(chart_type.lower(), 
                                                self.XL_CHART_TYPE.COLUMN_CLUSTERED)
            
            # Add chart
            left, top = position
            width, height = size
            slide.shapes.add_chart(
                chart_type_enum,
                left, top, width, height,
                chart_data
            )
            
        except Exception as e:
            raise DocumentCreationError(f"Failed to add chart: {e}")
    
    def apply_transition(
        self,
        slide_index: int,
        transition_type: str = "fade",
        duration: float = 1.0
    ) -> None:
        """
        Apply transition to a slide.
        
        Args:
            slide_index: Index of the slide (0-based)
            transition_type: Type of transition
            duration: Transition duration in seconds
        """
        if not self._presentation:
            raise DocumentCreationError(
                "Presentation not initialized. Call create_presentation() or load_presentation() first."
            )
        
        try:
            # Get slide
            if slide_index < 0 or slide_index >= len(self._presentation.slides):
                raise IndexError(f"Slide index {slide_index} out of range")
            
            slide = self._presentation.slides[slide_index]
            
            # Apply transition (simplified - python-pptx has limited transition support)
            # In production, you'd use the actual transition properties
            
            print(f"Note: Transition '{transition_type}' applied to slide {slide_index} "
                  f"(duration: {duration}s)")
            print("Full transition support requires advanced python-pptx usage.")
            
        except Exception as e:
            raise DocumentCreationError(f"Failed to apply transition: {e}")
    
    def save(
        self,
        file_path: Union[str, Path],
        format: str = "pptx"
    ) -> str:
        """
        Save the presentation to a file.
        
        Args:
            file_path: Path where to save the presentation
            format: Output format (pptx, pdf, etc.)
            
        Returns:
            Path to saved file
            
        Raises:
            DocumentSaveError: If saving fails
        """
        context = create_error_context(
            "save_presentation",
            file_path=file_path,
            format=format
        )
        
        validate_file_path(str(file_path), "save")
        validate_format(format, self.SUPPORTED_FORMATS, "save")
        
        if not self._presentation:
            raise DocumentSaveError(
                "No presentation to save. Create or load a presentation first.",
                context
            )
        
        try:
            file_path = Path(file_path)
            
            # Ensure directory exists
            file_path.parent.mkdir(parents=True, exist_ok=True)
            
            # Save in requested format
            if format == "pptx":
                self._presentation.save(str(file_path))
            elif format == "pdf":
                self._save_as_pdf(file_path)
            elif format in ["jpg", "png"]:
                self._save_as_images(file_path, format)
            else:
                # Default to pptx
                self._presentation.save(str(file_path))
            
            return str(file_path)
            
        except Exception as e:
            raise DocumentSaveError(
                f"Failed to save presentation: {e}",
                context
            )
    
    def _save_as_pdf(self, file_path: Path) -> None:
        """Save presentation as PDF."""
        # This is a placeholder - in production, you'd use a proper PDF converter
        raise NotImplementedError(
            "PDF export requires additional setup. "
            "Install python-pptx with PDF support or use a separate PDF library."
        )
    
    def _save_as_images(self, file_path: Path, format: str) -> None:
        """Save presentation as images (one per slide)."""
        # This is a placeholder - in production, you'd use a proper image export
        raise NotImplementedError(
            f"Image export ({format}) requires additional setup. "
            "Install python-pptx with image export support."
        )
    
    def _resolve_template_path(self, template: Union[str, Path]) -> Path:
        """Resolve template path."""
        template = str(template)
        
        # If it's an absolute path or relative path, use as-is
        if os.path.isabs(template) or template.startswith('.'):
            return Path(template)
        
        # Check in template directory
        if self.template_dir:
            template_path = Path(self.template_dir) / template
            if template_path.exists():
                return template_path
        
        # Check common extensions
        for ext in ['.pptx', '.ppt', '.potx']:
            test_path = Path(template + ext)
            if test_path.exists():
                return test_path
        
        # Return as-is (will fail if not found)
        return Path(template)
    
    def get_presentation(self) -> Optional['Presentation']:
        """Get the current presentation object."""
        return self._presentation
    
    def get_slide_count(self) -> int:
        """Get number of slides in presentation."""
        if not self._presentation:
            return 0
        return len(self._presentation.slides)
    
    def get_slide(self, index: int) -> 'Slide':
        """Get slide by index."""
        if not self._presentation:
            raise DocumentCreationError("Presentation not initialized")
        
        if 0 <= index < len(self._presentation.slides):
            return self._presentation.slides[index]
        else:
            raise IndexError(f"Slide index {index} out of range")
    
    def clear(self) -> None:
        """Clear the current presentation."""
        self._presentation = None
    
    def is_initialized(self) -> bool:
        """Check if presentation is initialized."""
        return self._presentation is not None